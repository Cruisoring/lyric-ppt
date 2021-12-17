from file_helper import all_lines
from pptx import Presentation
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_AUTO_SIZE
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR_INDEX as MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import ColorFormat, RGBColor

from xpinyin import Pinyin
import re


DEFAULT_TEMPLATE_FILE = 'templates\\template.pptx'
STOPS = '[，！。；：、「」，：]'


class Lyric:

    def __init__(self, lines, template_file: str = None) -> None:
        self.lines = lines
        if template_file is None:
            template_file = DEFAULT_TEMPLATE_FILE
        self.template = template_file
        self.ppt = Presentation(template_file)

        self.pinyin = Pinyin()

    def generate(self, withPinyin=True):
        start = -1
        for i in range(len(self.lines)):
            line = self.lines[i]
            if line.startswith('# '):
                title = line[2:].strip('\n \t')
                self.title(title)
            elif line.strip():
                if start < 0:
                    start = i
            else:
                if start > 0:
                    contents = self.lines[start:i]
                    self.new_content(contents, withPinyin)
                    start = -1

        self.ppt.save('Output.pptx')

    def title(self, title: str):
        slide = self.ppt.slides.add_slide(self.ppt.slide_layouts[0])
        title_shape = slide.shapes.title
        title_shape.text = title
        slide.placeholders[1].text = self.pinyin.get_pinyin(
            title, tone_marks='marks')

    def as_pinyin(self, line):
        py_text = self.pinyin.get_pinyin(line, tone_marks='marks')
        result = re.sub(r'\-[\s ]\-', ' ', py_text)
        return result

    def add_paragraph(self, frame, pt, lines):
        p = frame.add_paragraph()
        if len(lines) <= 2:
            lines.insert(0, ' ')
        elif len(lines) == 3:
            lines.insert(0, ' ')

        p.text = '\n'.join(lines)
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        p.font.size = Pt(pt)

    def title2(self, title: str):
        blank = self.ppt.slide_layouts[1]
        slide = self.ppt.slides.add_slide(blank)
        shapes = slide.shapes
        t_box = shapes.add_textbox(
            Inches(0.75), Inches(2.5), Inches(8.5), Inches(2))
        frame = t_box.text_frame

        self.add_paragraph(frame, 80, [title])
        self.add_paragraph(frame, 60, [self.as_pinyin(title)])

    def new_content(self, lines, withPinyin=True) -> any:
        blank = self.ppt.slide_layouts[1]
        slide = self.ppt.slides.add_slide(blank)
        shapes = slide.shapes
        t_box = shapes.add_textbox(
            Inches(0.25), Inches(0.15), Inches(12.5), Inches(7))
        frame = t_box.text_frame

        lines = [re.sub(STOPS, " ", l).strip(' \n') for l in lines]
        if withPinyin:
            max_len = max(len(l) for l in lines)
            pt = 56
            if len(lines) > 4:
                pt = 46
            elif max_len > 12:
                pt = 50
            elif len(lines) < 3:
                pt = 64

            self.add_paragraph(frame, pt, lines)

            lines = [self.as_pinyin(l) for l in lines]
            self.add_paragraph(frame, pt-16, lines)
        else:
            max_len = max(len(l) for l in lines)
            pt = 56
            if len(lines) > 8:
                pt = 46
            elif max_len > 12:
                pt = 50
            elif len(lines) < 6:
                pt = 64

            self.add_paragraph(frame, pt, lines)


def together(lyrics_file):
    lyric_list = [l.strip() for l in all_lines(lyrics_file) if l.strip()]
    lyric_lines = [all_lines(lyric) for lyric in lyric_list]
    all_lyric_lines = [l for lines in lyric_lines for l in lines]
    return all_lyric_lines


if __name__ == '__main__':
    lines = all_lines('平安夜')
    # lines = together('2021-12-19')
    l = Lyric(lines)
    l.generate(False)
