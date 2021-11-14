from pptx import Presentation
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_AUTO_SIZE
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR_INDEX as MSO_THEME_COLOR, MSO_FILL_TYPE
from pptx.dml.color import ColorFormat, RGBColor
from pptx.util import Inches, Pt

ONE_INCH = Inches(1.0)

ppt = Presentation()
# ['Title Slide', 'Title and Content', 'Section Header', 'Two Content', 'Comparison', 'Title Only', 'Blank', 'Content with Caption', 'Picture with Caption', 'Title and Vertical Text', 'Vertical Title and Text']

def print_inches(shape):
    left = shape.left / ONE_INCH
    top = shape.top / ONE_INCH
    width = shape.width / ONE_INCH
    height = shape.height / ONE_INCH
    print(f'Inches(L={left}, T={top}, W={width}, H={height})')

title_layout = ppt.slide_layouts[0]
# fill = title_layout.background.fill

# #fill.fore_color.rgb = RGBColor(255, 0, 0)

# title_slide = p.slides.add_slide(title_layout)

for shape in title_layout.shapes:
    print_inches(shape)



# blank = p.slide_layouts[6]
# print(blank.name)

# slide = p.slides.add_slide(blank)

blank = ppt.slide_layouts[6]
slide = ppt.slides.add_slide(blank)
shapes = slide.shapes
t_box = shapes.add_textbox(Inches(0.25), Inches(0.25), Inches(9.5), Inches(7))
frame = t_box.text_frame

p = frame.add_paragraph()
p.text = 'This is the body\nLine 2\nLine 3'
p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
p.font.size = Pt(60)

p = frame.add_paragraph()
p.text = 'pin-yin bieiieieie'
p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
p.font.size = Pt(40)

ppt.save('Output.pptx')
